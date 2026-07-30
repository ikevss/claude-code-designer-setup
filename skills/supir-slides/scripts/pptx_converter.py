#!/usr/bin/env python3
"""PPTX 转换器 - 解析 PPTX 并生成 open-slide 兼容的 React 组件"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("错误: 需要安装 python-pptx: pip install python-pptx", file=sys.stderr)
    sys.exit(2)


def extract_text_from_shape(shape):
    """从 shape 中提取文本"""
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""


def extract_table_data(shape):
    """从表格 shape 中提取数据"""
    if not shape.has_table:
        return None

    table = shape.table
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    return rows


def shape_to_css(shape, slide_width, slide_height):
    """将 shape 坐标转换为 CSS 样式"""
    left_pct = (shape.left / slide_width) * 100
    top_pct = (shape.top / slide_height) * 100
    width_pct = (shape.width / slide_width) * 100
    height_pct = (shape.height / slide_height) * 100

    return {
        "position": "absolute",
        "left": f"{left_pct:.1f}%",
        "top": f"{top_pct:.1f}%",
        "width": f"{width_pct:.1f}%",
        "height": f"{height_pct:.1f}%",
    }


def parse_pptx(pptx_path):
    """解析 PPTX 文件"""
    prs = Presentation(pptx_path)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    result = {
        "file": pptx_path,
        "dimensions": {
            "width_emu": slide_width,
            "height_emu": slide_height,
            "width_px": 1920,
            "height_px": 1080,
        },
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "index": i + 1,
            "layout": slide.slide_layout.name if slide.slide_layout else "unknown",
            "elements": []
        }

        for shape in slide.shapes:
            element = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "css": shape_to_css(shape, slide_width, slide_height)
            }

            text = extract_text_from_shape(shape)
            if text:
                element["text"] = text

            table_data = extract_table_data(shape)
            if table_data:
                element["table"] = table_data

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    element["image"] = {
                        "content_type": image.content_type,
                        "size": len(image.blob)
                    }
                except:
                    pass

            slide_data["elements"].append(element)

        result["slides"].append(slide_data)

    return result


def generate_react_component(slide_data, output_dir):
    """为单个幻灯片生成 React 组件"""
    index = slide_data["index"]
    filename = f"{index:02d}-slide.tsx"

    elements_jsx = []
    for elem in slide_data["elements"]:
        css_str = json.dumps(elem["css"])

        if "text" in elem:
            text = elem["text"].replace('"', '\\"').replace('\n', '\\n')
            elements_jsx.append(f'''
      <div style={{{{ ...{css_str}, display: "flex", alignItems: "center", justifyContent: "center" }}}}>
        <p style={{{ fontSize: "1.2rem", margin: 0 }}}>{text}</p>
      </div>''')
        elif "table" in elem:
            table_html = "<table><tbody>"
            for row in elem["table"]:
                table_html += "<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>"
            table_html += "</tbody></table>"
            elements_jsx.append(f'''
      <div style={{{{ ...{css_str}, overflow: "auto" }}}}>
        {table_html}
      </div>''')

    component = f'''export default function Slide{index:02d}() {{
  return (
    <div style={{{{ position: "relative", width: "1920px", height: "1080px", overflow: "hidden" }}}}>
      {"".join(elements_jsx)}
    </div>
  );
}}
'''
    output_path = Path(output_dir) / filename
    output_path.write_text(component, encoding="utf-8")
    return str(output_path)


def main():
    parser = argparse.ArgumentParser(description="PPTX 转 open-slide React 组件")
    parser.add_argument("--pptx", required=True, help="PPTX 文件路径")
    parser.add_argument("--output-dir", default="./slides", help="输出目录")
    parser.add_argument("--json", action="store_true", help="输出 JSON 格式")
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"错误: 文件不存在 {args.pptx}", file=sys.stderr)
        return 1

    try:
        result = parse_pptx(args.pptx)
    except Exception as e:
        print(f"解析失败: {e}", file=sys.stderr)
        return 1

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0

    os.makedirs(args.output_dir, exist_ok=True)
    generated = []
    for slide_data in result["slides"]:
        path = generate_react_component(slide_data, args.output_dir)
        generated.append(path)

    print(json.dumps({
        "status": "success",
        "slides_count": len(generated),
        "output_dir": args.output_dir,
        "files": generated
    }))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
