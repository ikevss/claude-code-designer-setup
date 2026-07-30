#!/usr/bin/env python3
"""图片转可编辑 PPT - 使用 PaddleOCR-VL API"""

import argparse
import json
import os
import sys
import time
import urllib.request
import urllib.error
from pathlib import Path

# 从 .env 文件读取配置
def load_env():
    """从 .env 文件加载环境变量"""
    env_path = Path(__file__).parent.parent / ".env"
    if env_path.exists():
        with open(env_path, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line and not line.startswith("#") and "=" in line:
                    key, value = line.split("=", 1)
                    os.environ.setdefault(key.strip(), value.strip())

load_env()

# API 配置
JOB_URL = "https://paddleocr.aistudio-app.com/api/v2/ocr/jobs"
DEFAULT_TOKEN = os.environ.get("PADDLE_OCR_TOKEN", "")
MODEL = "PaddleOCR-VL-1.6"


def submit_ocr_job(file_path, token):
    """提交 OCR 作业"""
    headers = {
        "Authorization": f"bearer {token}"
    }

    if file_path.startswith("http"):
        # URL 模式
        headers["Content-Type"] = "application/json"
        payload = {
            "fileUrl": file_path,
            "model": MODEL,
            "optionalPayload": {
                "useDocOrientationClassify": False,
                "useDocUnwarping": False,
                "useChartRecognition": False
            }
        }
        req = urllib.request.Request(
            JOB_URL,
            data=json.dumps(payload).encode("utf-8"),
            headers=headers,
            method="POST"
        )
    else:
        # 本地文件模式
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        boundary = "----WebKitFormBoundary7MA4YWxkTrZu0gW"
        file_name = os.path.basename(file_path)

        with open(file_path, "rb") as f:
            file_data = f.read()

        body = (
            f"--{boundary}\r\n"
            f'Content-Disposition: form-data; name="model"\r\n\r\n'
            f"{MODEL}\r\n"
            f"--{boundary}\r\n"
            f'Content-Disposition: form-data; name="optionalPayload"\r\n\r\n'
            f"{json.dumps({'useDocOrientationClassify': False, 'useDocUnwarping': False, 'useChartRecognition': False})}\r\n"
            f"--{boundary}\r\n"
            f'Content-Disposition: form-data; name="file"; filename="{file_name}"\r\n'
            f"Content-Type: application/octet-stream\r\n\r\n"
        ).encode("utf-8") + file_data + f"\r\n--{boundary}--\r\n".encode("utf-8")

        headers["Content-Type"] = f"multipart/form-data; boundary={boundary}"
        req = urllib.request.Request(JOB_URL, data=body, headers=headers, method="POST")

    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", "replace")
        raise RuntimeError(f"提交作业失败: {e.code} - {body}") from e

    if result.get("code") != 200:
        raise RuntimeError(f"提交作业失败: {result}")

    return result["data"]["jobId"]


def poll_job_status(job_id, token, max_wait=300):
    """轮询作业状态"""
    headers = {"Authorization": f"bearer {token}"}
    start_time = time.time()

    while True:
        if time.time() - start_time > max_wait:
            raise RuntimeError(f"作业超时（超过 {max_wait} 秒）")

        url = f"{JOB_URL}/{job_id}"
        req = urllib.request.Request(url, headers=headers)

        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                result = json.loads(resp.read().decode("utf-8"))
        except Exception as e:
            print(f"⚠️ 轮询失败，重试中: {e}", file=sys.stderr)
            time.sleep(5)
            continue

        data = result.get("data", {})
        state = data.get("state", "unknown")

        if state == "pending":
            print("⏳ 作业排队中...", file=sys.stderr)
        elif state == "running":
            progress = data.get("extractProgress", {})
            total = progress.get("totalPages", "?")
            extracted = progress.get("extractedPages", "?")
            print(f"🔄 识别中: {extracted}/{total} 页", file=sys.stderr)
        elif state == "done":
            print("✅ 识别完成", file=sys.stderr)
            return data.get("resultUrl", {}).get("jsonUrl")
        elif state == "failed":
            error_msg = data.get("errorMsg", "未知错误")
            raise RuntimeError(f"OCR 任务失败: {error_msg}")
        else:
            print(f"⚠️ 未知状态: {state}", file=sys.stderr)

        time.sleep(5)


def fetch_results(jsonl_url):
    """获取 OCR 结果"""
    req = urllib.request.Request(jsonl_url)

    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            content = resp.read().decode("utf-8")
    except Exception as e:
        raise RuntimeError(f"获取结果失败: {e}") from e

    results = []
    for line in content.strip().split("\n"):
        line = line.strip()
        if not line:
            continue
        try:
            data = json.loads(line)
            result = data.get("result", {})
            for page_result in result.get("layoutParsingResults", []):
                markdown = page_result.get("markdown", {})
                text = markdown.get("text", "")
                if text:
                    results.append(text)
        except json.JSONDecodeError:
            continue

    return results


def create_pptx_from_text(text_blocks, output_path):
    """根据识别结果创建可编辑 PPTX"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise RuntimeError("需要安装 python-pptx: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    for page_text in text_blocks:
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 按行分割文本
        lines = [line.strip() for line in page_text.split("\n") if line.strip()]

        y_offset = Inches(0.8)
        for line in lines[:20]:  # 最多 20 行
            txBox = slide.shapes.add_textbox(
                Inches(1), y_offset, Inches(11), Inches(0.4)
            )
            tf = txBox.text_frame
            tf.text = line
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
            y_offset += Inches(0.32)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="图片转可编辑 PPT")
    parser.add_argument("--input", required=True, help="输入图片路径或 URL")
    parser.add_argument("--output", default="output.pptx", help="输出 PPTX 路径")
    parser.add_argument("--token", default="", help="PaddleOCR Token")
    parser.add_argument("--confirm", action="store_true", help="跳过确认直接执行")
    args = parser.parse_args()

    token = args.token or DEFAULT_TOKEN
    if not token:
        print("❌ 错误: 未配置 PaddleOCR Token", file=sys.stderr)
        print("请在 .env 文件中设置 PADDLE_OCR_TOKEN", file=sys.stderr)
        return 1

    input_path = args.input
    is_url = input_path.startswith("http")

    if not is_url:
        input_path = Path(input_path)
        if not input_path.exists():
            print(f"❌ 错误: 文件不存在 {input_path}", file=sys.stderr)
            return 1
        file_size_mb = input_path.stat().st_size / (1024 * 1024)
        input_str = str(input_path)
    else:
        file_size_mb = 0
        input_str = input_path

    # 用户确认提示
    if not args.confirm:
        print("=" * 60)
        print("⚠️  图片转可编辑 PPT 需要消耗较多 Token 和时间")
        print("=" * 60)
        print(f"输入文件: {input_str}")
        if not is_url:
            print(f"文件大小: {file_size_mb:.2f} MB")
        print(f"预计耗时: 60-180 秒（取决于图片复杂度）")
        print("=" * 60)
        print("\n是否继续？(y/n)")
        response = input().strip().lower()
        if response != "y":
            print("已取消")
            return 0

    # 提交 OCR 作业
    print("📤 正在提交识别任务...")
    try:
        job_id = submit_ocr_job(input_str, token)
        print(f"✅ 任务已提交，ID: {job_id}")
    except Exception as e:
        print(f"❌ 提交失败: {e}", file=sys.stderr)
        return 1

    # 轮询状态
    print("🔍 正在识别图片文字...")
    try:
        jsonl_url = poll_job_status(job_id, token)
    except Exception as e:
        print(f"❌ 识别失败: {e}", file=sys.stderr)
        return 1

    # 获取结果
    print("📥 正在获取识别结果...")
    try:
        text_blocks = fetch_results(jsonl_url)
        print(f"✅ 识别完成，共 {len(text_blocks)} 页")
    except Exception as e:
        print(f"❌ 获取结果失败: {e}", file=sys.stderr)
        return 1

    # 创建 PPTX
    print("📝 正在生成可编辑 PPT...")
    try:
        output_path = create_pptx_from_text(text_blocks, args.output)
        print(f"✅ PPTX 已生成: {output_path}")
    except Exception as e:
        print(f"❌ 生成失败: {e}", file=sys.stderr)
        return 1

    # 输出结果
    result = {
        "status": "success",
        "input": input_str,
        "output": output_path,
        "pages": len(text_blocks)
    }
    print(json.dumps(result))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
