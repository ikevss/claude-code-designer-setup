#!/usr/bin/env python3
"""pptx-replicate / parse_ooxml：解析参考 PPTX 的 OOXML 结构，抠出复刻所需的精确数据。

python-pptx 直接读 OOXML，能拿到图片给不了的「精确结构化数据」：画布尺寸、每页版式、
每个 shape 的坐标(left/top/width/height)、文本 run 的字号/字体/颜色/粗体、表格单元格内容、
layout/master 继承形状、图片 hash/尺寸/裁剪信息和基础设计摘要。
这些是从 PNG 里反推不出来的，复刻一份 PPT 时必须从这里抠。

与 render_png.py 互补：render_png 让你「看图」学版式风格，parse_ooxml 给你「精确数值」去复刻。

默认打印人类可读的逐页大纲（版式 + 文本 + 坐标摘要）；加 --json 输出完整结构化数据
（可直接序列化喂给 paipai-slides）。

用法::

    python3 scripts/parse_ooxml.py --pptx-path "tmp/upload/极米科技深度报告@202203.pptx"
    python3 scripts/parse_ooxml.py --pptx-path "tmp/upload/x.pptx" --json > structure.json
    python3 scripts/parse_ooxml.py --pptx-path "tmp/upload/x.pptx" --slides 1-5
    python3 scripts/parse_ooxml.py --pptx-path "tmp/upload/x.pptx" --json --extract-images work01/assets

依赖 python-pptx（pod 内已预装）。失败时往 stderr 打印错误并以非零退出码退出。
"""

import argparse
import hashlib
import json
import sys
from collections import Counter
from pathlib import Path

EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def _emu_to_in(v):
    return round(v / EMU_PER_INCH, 3) if v is not None else None


def _emu_to_pt(v):
    return round(v / EMU_PER_PT, 1) if v is not None else None


def _maybe(v):
    return v if v not in ("", None, [], {}) else None


def _enum_str(v):
    if v is None:
        return None
    try:
        return str(v)
    except Exception:
        return repr(v)


def _rgb_of(color):
    """尽力取颜色：优先 RGB，其次主题色名；取不到返回 None。"""
    try:
        if color is None:
            return None
        # color.type 为 None 时未显式着色
        if color.type is None:
            return None
        try:
            return "#" + str(color.rgb)
        except Exception:
            pass
        try:
            theme = color.theme_color
            if theme is not None:
                return f"theme:{theme}"
        except Exception:
            pass
    except Exception:
        pass
    return None


def _fill_format_info(fill):
    info = {}
    try:
        if fill.type is not None:
            info["type"] = _enum_str(fill.type)
    except Exception:
        pass
    try:
        rgb = _rgb_of(fill.fore_color)
        if rgb:
            info["fore_color"] = rgb
    except Exception:
        pass
    try:
        rgb = _rgb_of(fill.back_color)
        if rgb:
            info["back_color"] = rgb
    except Exception:
        pass
    try:
        if fill.transparency:
            info["transparency"] = fill.transparency
    except Exception:
        pass
    return _maybe(info)


def _fill_info(shape):
    try:
        return _fill_format_info(shape.fill)
    except Exception:
        return None


def _line_info(shape):
    try:
        line = shape.line
    except Exception:
        return None
    info = {}
    try:
        rgb = _rgb_of(line.color)
        if rgb:
            info["color"] = rgb
    except Exception:
        pass
    try:
        if line.width is not None:
            info["width_pt"] = _emu_to_pt(int(line.width))
    except Exception:
        pass
    try:
        if line.dash_style is not None:
            info["dash_style"] = _enum_str(line.dash_style)
    except Exception:
        pass
    try:
        if line.transparency:
            info["transparency"] = line.transparency
    except Exception:
        pass
    return _maybe(info)


def _font_info(font):
    """提炼 run/字体信息，只保留非空字段。"""
    info = {}
    try:
        if font.name:
            info["font"] = font.name
    except Exception:
        pass
    try:
        if font.size is not None:
            info["size_pt"] = _emu_to_pt(int(font.size))
    except Exception:
        pass
    try:
        if font.bold:
            info["bold"] = True
    except Exception:
        pass
    try:
        if font.italic:
            info["italic"] = True
    except Exception:
        pass
    try:
        if font.underline:
            info["underline"] = _enum_str(font.underline)
    except Exception:
        pass
    rgb = _rgb_of(getattr(font, "color", None))
    if rgb:
        info["color"] = rgb
    return info


def _text_frame_info(text_frame):
    info = {}
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            value = getattr(text_frame, attr)
            if value is not None:
                info[attr + "_pt"] = _emu_to_pt(int(value))
        except Exception:
            pass
    try:
        if text_frame.vertical_anchor is not None:
            info["vertical_anchor"] = _enum_str(text_frame.vertical_anchor)
    except Exception:
        pass
    try:
        if text_frame.word_wrap is not None:
            info["word_wrap"] = text_frame.word_wrap
    except Exception:
        pass
    try:
        if text_frame.auto_size is not None:
            info["auto_size"] = _enum_str(text_frame.auto_size)
    except Exception:
        pass
    return _maybe(info)


def _paragraphs(text_frame):
    """逐段抽取：合并整段文字，并保留首个非空 run 的字体特征作代表。"""
    out = []
    for p in text_frame.paragraphs:
        text = "".join(r.text for r in p.runs) or p.text or ""
        text = text.strip()
        if not text:
            continue
        para = {"text": text}
        try:
            if p.level:
                para["level"] = p.level
        except Exception:
            pass
        try:
            if p.alignment is not None:
                para["align"] = str(p.alignment)
        except Exception:
            pass
        try:
            p_font = _font_info(p.font)
            if p_font:
                para["paragraph_font"] = p_font
        except Exception:
            pass
        runs = []
        # 取首个有内容的 run 的字体作整段代表（深度报告里多数整段同风格）
        for r in p.runs:
            text = r.text.strip()
            if not text:
                continue
            fi = _font_info(r.font)
            if "font" not in para and fi:
                para["font"] = fi
            run = {"text": text}
            if fi:
                run["font"] = fi
            runs.append(run)
        if runs:
            para["runs"] = runs
        out.append(para)
    return out


def _cell_detail(cell):
    detail = {"text": cell.text.strip()}
    fill = _fill_format_info(cell.fill)
    if fill:
        detail["fill"] = fill
    try:
        paras = _paragraphs(cell.text_frame)
        if paras:
            detail["paragraphs"] = paras
    except Exception:
        pass
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            value = getattr(cell, attr)
            if value is not None:
                detail[attr + "_pt"] = _emu_to_pt(int(value))
        except Exception:
            pass
    for attr in ("is_merge_origin", "is_spanned", "span_height", "span_width"):
        try:
            value = getattr(cell, attr)
            if attr.startswith("span_") and value == 1:
                continue
            if value:
                detail[attr] = value
        except Exception:
            pass
    return detail


def _table(shape):
    tbl = shape.table
    rows = []
    cell_details = []
    for row in tbl.rows:
        cells = []
        details = []
        for cell in row.cells:
            cells.append(cell.text.strip())
            details.append(_cell_detail(cell))
        rows.append(cells)
        cell_details.append(details)
    out = {
        "rows": len(tbl.rows._tbl.tr_lst),
        "cols": len(rows[0]) if rows else 0,
        "cells": rows,
    }
    try:
        out["row_heights_in"] = [_emu_to_in(r.height) for r in tbl.rows]
    except Exception:
        pass
    try:
        out["col_widths_in"] = [_emu_to_in(c.width) for c in tbl.columns]
    except Exception:
        pass
    if cell_details:
        out["cell_details"] = cell_details
    return out


def _chart(shape):
    chart = shape.chart
    out = {"type": _enum_str(chart.chart_type)}
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                out["title"] = title
    except Exception:
        pass
    try:
        if chart.has_legend:
            out["legend"] = {
                "position": _enum_str(chart.legend.position),
                "include_in_layout": chart.legend.include_in_layout,
            }
    except Exception:
        pass
    series = []
    try:
        for s in chart.series:
            item = {}
            try:
                item["name"] = str(s.name)
            except Exception:
                pass
            for attr in ("values", "x_values", "y_values"):
                try:
                    values = getattr(s, attr)
                    if values is not None:
                        item[attr] = list(values)
                except Exception:
                    pass
            if item:
                series.append(item)
    except Exception:
        pass
    if series:
        out["series"] = series
    return out


def _picture_info(shape, image_dir=None, slide_index=None, shape_index=None, source="slide"):
    out = {}
    try:
        image = shape.image
    except Exception:
        return {"present": True}

    blob = None
    try:
        blob = image.blob
        out["sha256"] = hashlib.sha256(blob).hexdigest()
        out["bytes"] = len(blob)
    except Exception:
        pass
    for attr in ("filename", "content_type", "ext"):
        try:
            value = getattr(image, attr)
            if value:
                out[attr] = value
        except Exception:
            pass
    try:
        w_px, h_px = image.size
        out["image_px"] = {"width": w_px, "height": h_px}
    except Exception:
        pass
    try:
        out["rId"] = shape._element.blip_rId
    except Exception:
        pass
    crop = {}
    for attr in ("crop_left", "crop_right", "crop_top", "crop_bottom"):
        try:
            value = getattr(shape, attr)
            if value:
                crop[attr] = round(float(value), 5)
        except Exception:
            pass
    if crop:
        out["crop"] = crop

    if image_dir and blob is not None:
        try:
            image_dir.mkdir(parents=True, exist_ok=True)
            suffix = out.get("ext") or "img"
            sha = out.get("sha256")
            if sha:
                filename = f"image_{sha[:12]}.{suffix}"
            else:
                filename = f"s{slide_index:03d}_{source}_shape{shape_index:03d}.{suffix}"
            path = image_dir / filename
            if not path.exists():
                path.write_bytes(blob)
            out["asset_path"] = str(path)
        except Exception as exc:
            out["asset_error"] = str(exc)
    return out


def _placeholder_info(shape):
    try:
        if not shape.is_placeholder:
            return None
        ph = shape.placeholder_format
    except Exception:
        return None
    info = {}
    try:
        info["idx"] = ph.idx
    except Exception:
        pass
    try:
        info["type"] = _enum_str(ph.type)
    except Exception:
        pass
    return _maybe(info)


def _shape_dict(shape, source="slide", image_dir=None, slide_index=None, shape_index=0):
    """单个 shape → 结构化 dict；递归处理 group。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    d = {
        "name": shape.name,
        "type": str(shape.shape_type),
        "source": source,
        "shape_index": shape_index,
        "pos_in": {
            "left": _emu_to_in(shape.left),
            "top": _emu_to_in(shape.top),
            "width": _emu_to_in(shape.width),
            "height": _emu_to_in(shape.height),
        },
    }
    try:
        d["shape_id"] = shape.shape_id
    except Exception:
        pass
    try:
        if shape.auto_shape_type is not None:
            d["auto_shape_type"] = _enum_str(shape.auto_shape_type)
    except Exception:
        pass
    placeholder = _placeholder_info(shape)
    if placeholder:
        d["placeholder"] = placeholder
    fill = _fill_info(shape)
    if fill:
        d["fill"] = fill
    line = _line_info(shape)
    if line:
        d["line"] = line

    # group：递归
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            d["children"] = [
                _shape_dict(c, source, image_dir, slide_index, child_index)
                for child_index, c in enumerate(shape.shapes, start=1)
            ]
            return d
    except Exception:
        pass

    # 表格
    try:
        if shape.has_table:
            d["table"] = _table(shape)
            return d
    except Exception:
        pass

    # 图表
    try:
        if shape.has_chart:
            d["chart"] = _chart(shape)
            return d
    except Exception:
        pass

    # 图片
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            d["picture"] = _picture_info(shape, image_dir, slide_index, shape_index, source)
    except Exception:
        pass

    # 文本
    try:
        if shape.has_text_frame:
            tf = _text_frame_info(shape.text_frame)
            if tf:
                d["text_frame"] = tf
            paras = _paragraphs(shape.text_frame)
            if paras:
                d["paragraphs"] = paras
    except Exception:
        pass

    return d


def _shape_list(shapes, source, image_dir=None, slide_index=None):
    return [
        _shape_dict(s, source, image_dir, slide_index, shape_index)
        for shape_index, s in enumerate(shapes, start=1)
    ]


def _iter_shape_dicts(shapes):
    for shape in shapes:
        yield shape
        for child in shape.get("children", []):
            yield from _iter_shape_dicts([child])


def _collect_design_summary(result):
    fonts = Counter()
    font_sizes = Counter()
    font_signatures = Counter()
    colors = Counter()
    counts = Counter()
    large_text_runs = []

    def remember_color(value):
        if value:
            colors[value] += 1

    def remember_font(font, text="", slide_index=None, shape_name="", source=""):
        if not font:
            return
        name = font.get("font")
        size = font.get("size_pt")
        if name:
            fonts[name] += 1
        if size is not None:
            font_sizes[str(size)] += 1
        remember_color(font.get("color"))
        if name or size is not None or font.get("bold") or font.get("italic"):
            signature = [
                name or "unknown-font",
                f"{size}pt" if size is not None else "unknown-size",
                "bold" if font.get("bold") else "regular",
            ]
            if font.get("italic"):
                signature.append("italic")
            if font.get("color"):
                signature.append(font["color"])
            font_signatures[" | ".join(signature)] += 1
        if text and size is not None:
            large_text_runs.append(
                {
                    "text": text[:80],
                    "size_pt": size,
                    "font": name,
                    "bold": bool(font.get("bold")),
                    "color": font.get("color"),
                    "slide": slide_index,
                    "source": source,
                    "shape": shape_name,
                }
            )

    for slide in result["slides"]:
        shape_groups = [slide.get("shapes", [])]
        inherited = slide.get("inherited_shapes") or {}
        shape_groups.extend(inherited.get(k, []) for k in ("layout", "master"))
        for group in shape_groups:
            for sh in _iter_shape_dicts(group):
                counts["shapes"] += 1
                if sh.get("picture"):
                    counts["pictures"] += 1
                if sh.get("table"):
                    counts["tables"] += 1
                if sh.get("chart"):
                    counts["charts"] += 1
                fill = sh.get("fill") or {}
                remember_color(fill.get("fore_color"))
                remember_color(fill.get("back_color"))
                line = sh.get("line") or {}
                remember_color(line.get("color"))
                for para in sh.get("paragraphs", []):
                    for font in (para.get("font"), para.get("paragraph_font")):
                        remember_font(
                            font,
                            para.get("text", ""),
                            slide.get("index"),
                            sh.get("name", ""),
                            sh.get("source", ""),
                        )
                    for run in para.get("runs", []):
                        remember_font(
                            run.get("font") or {},
                            run.get("text", ""),
                            slide.get("index"),
                            sh.get("name", ""),
                            sh.get("source", ""),
                        )
                table = sh.get("table") or {}
                for row in table.get("cell_details", []):
                    for cell in row:
                        fill = cell.get("fill") or {}
                        remember_color(fill.get("fore_color"))
                        for para in cell.get("paragraphs", []):
                            remember_font(
                                para.get("font") or {},
                                para.get("text", ""),
                                slide.get("index"),
                                sh.get("name", ""),
                                sh.get("source", ""),
                            )

    large_text_runs = sorted(
        large_text_runs,
        key=lambda item: (item.get("size_pt") or 0, len(item.get("text") or "")),
        reverse=True,
    )[:20]
    return {
        "top_fonts": [{"value": k, "count": v} for k, v in fonts.most_common(12)],
        "top_font_sizes": [
            {"value": k, "count": v} for k, v in font_sizes.most_common(16)
        ],
        "top_font_signatures": [
            {"value": k, "count": v} for k, v in font_signatures.most_common(20)
        ],
        "largest_text_runs": large_text_runs,
        "top_colors": [{"value": k, "count": v} for k, v in colors.most_common(20)],
        "counts": dict(counts),
    }


def parse(pptx_path, include_inherited=True, image_dir=None):
    from pptx import Presentation

    prs = Presentation(pptx_path)
    w, h = prs.slide_width, prs.slide_height
    result = {
        "file": pptx_path,
        "dimensions": {
            "width_emu": w,
            "height_emu": h,
            "width_in": _emu_to_in(w),
            "height_in": _emu_to_in(h),
            "ratio": round(w / h, 4) if h else None,
        },
        "slide_count": len(prs.slides),
        "slides": [],
    }
    for i, slide in enumerate(prs.slides):
        layout = slide.slide_layout.name if slide.slide_layout else None
        slide_index = i + 1
        shapes = _shape_list(slide.shapes, "slide", image_dir, slide_index)
        slide_dict = {"index": slide_index, "layout": layout, "shapes": shapes}
        if include_inherited:
            inherited = {}
            try:
                inherited["layout"] = _shape_list(
                    slide.slide_layout.shapes, "layout", image_dir, slide_index
                )
            except Exception:
                inherited["layout"] = []
            try:
                inherited["master"] = _shape_list(
                    slide.slide_layout.slide_master.shapes,
                    "master",
                    image_dir,
                    slide_index,
                )
            except Exception:
                inherited["master"] = []
            slide_dict["inherited_shapes"] = inherited
        result["slides"].append(slide_dict)
    result["design_summary"] = _collect_design_summary(result)
    return result


def _parse_slide_range(spec, total):
    """'1-5' / '3' / '2,4,6' → 1-based 索引集合。"""
    if not spec:
        return set(range(1, total + 1))
    picked = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            picked.update(range(int(a), int(b) + 1))
        else:
            picked.add(int(part))
    return {n for n in picked if 1 <= n <= total}


def _print_human(result, wanted):
    dim = result["dimensions"]
    print(
        f"画布: {dim['width_in']}×{dim['height_in']} in "
        f"(ratio={dim['ratio']}, {dim['width_emu']}×{dim['height_emu']} EMU)"
    )
    print(f"总页数: {result['slide_count']}\n")
    for s in result["slides"]:
        if s["index"] not in wanted:
            continue
        print(f"=== S{s['index']} [layout: {s['layout']}] ===")
        _print_shapes(s["shapes"], indent=1)
        inherited = s.get("inherited_shapes") or {}
        for source in ("layout", "master"):
            shapes = inherited.get(source) or []
            if not shapes:
                continue
            print(f"  --- inherited {source} shapes ({len(shapes)}) ---")
            _print_shapes(shapes, indent=2)
        print()


def _print_shapes(shapes, indent):
    pad = "  " * indent
    for sh in shapes:
        p = sh["pos_in"]
        pos = f"@({p['left']},{p['top']}) {p['width']}×{p['height']}in"
        tag = sh["type"].split(".")[-1].split(" ")[0]
        line = f"{pad}• [{tag}] {sh['name']} {pos}"
        print(line)
        if "table" in sh:
            t = sh["table"]
            print(f"{pad}    表格 {t['rows']}×{t['cols']}:")
            for row in t["cells"]:
                print(f"{pad}    | " + " | ".join(c[:20] for c in row))
        if "chart" in sh:
            chart = sh["chart"]
            chart_type = chart.get("type") if isinstance(chart, dict) else chart
            print(f"{pad}    图表: {chart_type}")
        picture = sh.get("picture")
        if picture:
            if isinstance(picture, dict):
                sha = picture.get("sha256", "")[:12]
                size = picture.get("image_px")
                asset = picture.get("asset_path")
                bits = ["图片"]
                if sha:
                    bits.append(f"sha256={sha}")
                if size:
                    bits.append(f"{size['width']}×{size['height']}px")
                if asset:
                    bits.append(f"asset={asset}")
                print(f"{pad}    [" + ", ".join(bits) + "]")
            else:
                print(f"{pad}    [图片]")
        for para in sh.get("paragraphs", []):
            f = para.get("font", {})
            meta = []
            if f.get("size_pt"):
                meta.append(f"{f['size_pt']}pt")
            if f.get("font"):
                meta.append(f["font"])
            if f.get("bold"):
                meta.append("bold")
            if f.get("color"):
                meta.append(f["color"])
            lvl = ("  " * para.get("level", 0))
            meta_str = f"  <{', '.join(meta)}>" if meta else ""
            print(f"{pad}    {lvl}- {para['text'][:70]}{meta_str}")
        if "children" in sh:
            print(f"{pad}    (group, {len(sh['children'])} 子元素)")
            _print_shapes(sh["children"], indent + 2)


def main(argv=None):
    parser = argparse.ArgumentParser(
        prog="parse_ooxml",
        description="解析参考 PPTX 的 OOXML 结构（坐标/字号/颜色/表格），供复刻用。",
    )
    parser.add_argument(
        "--pptx-path", required=True, help="参考 PPTX 路径（相对工作区）"
    )
    parser.add_argument(
        "--json", action="store_true", help="输出完整结构化 JSON（默认输出人类可读大纲）"
    )
    parser.add_argument(
        "--slides", default="", help="只看指定页，如 '1-5' / '3' / '2,4,6'（默认全部）"
    )
    parser.add_argument(
        "--no-inherited",
        action="store_true",
        help="不解析 slide layout / master 中继承的 shape（默认解析，便于发现 Logo/页眉/页脚）",
    )
    parser.add_argument(
        "--extract-images",
        default="",
        help="可选：把 PPTX 中图片资产导出到指定目录，并在 JSON 中写入 asset_path",
    )
    args = parser.parse_args(argv)

    try:
        image_dir = Path(args.extract_images) if args.extract_images else None
        result = parse(
            args.pptx_path,
            include_inherited=not args.no_inherited,
            image_dir=image_dir,
        )
    except ImportError:
        print("缺少 python-pptx：请 pip install python-pptx（pod 内通常已预装）。",
              file=sys.stderr)
        return 2
    except Exception as e:
        print(f"解析失败: {e}", file=sys.stderr)
        return 1

    wanted = _parse_slide_range(args.slides, result["slide_count"])

    if args.json:
        if args.slides:
            result["slides"] = [s for s in result["slides"] if s["index"] in wanted]
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        _print_human(result, wanted)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
